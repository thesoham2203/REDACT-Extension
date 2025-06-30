import spacy
import os
import re
from faker import Faker
from pptx import Presentation

class PresentationRedactor:
    def __init__(self, spacy_model="en_core_web_sm"):
        self.fake = Faker()
        self.nlp = spacy.load(spacy_model)

        # Mapping redaction levels to sensitive types
        self.redaction_map = {
            0: [],
            25: ["ORG", "EMAIL", "PHONE"],
            50: ["ORG", "EMAIL", "PHONE", "MONEY", "IP"],
            75: ["ORG", "EMAIL", "PHONE", "MONEY", "IP", "DATE", "TIME"],
            100: ["ORG", "EMAIL", "PHONE", "MONEY", "IP", "DATE", "TIME", "ADDRESS", "PERSON"]
        }

    def extract_text_from_presentation(self, ppt_path):
        prs = Presentation(ppt_path)
        text_blocks = []
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_blocks.append({"text": shape.text})
                    text += shape.text + "\n"

        return text, text_blocks

    def detect_sensitive_data(self, text, redaction_labels):
        doc = self.nlp(text)
        sensitive_data = []

        for ent in doc.ents:
            if ent.label_ in redaction_labels:
                sensitive_data.append((ent.text, ent.label_))

        # Additional regex-based detection
        if "MONEY" in redaction_labels:
            money_matches = re.findall(r"\$\d+(\.\d{2})?", text)
            sensitive_data.extend([(match, "MONEY") for match in money_matches])
        if "TIME" in redaction_labels:
            time_matches = re.findall(r"\b\d{1,2}:\d{2}(?:\s?[APap][Mm])?\b", text)
            sensitive_data.extend([(match, "TIME") for match in time_matches])
        if "IP" in redaction_labels:
            ip_matches = re.findall(r"\b(?:\d{1,3}\.){3}\d{1,3}\b", text)
            sensitive_data.extend([(match, "IP") for match in ip_matches])

        return sensitive_data

    def generate_synthetic_data(self, label):
        if label == "PERSON":
            return self.fake.name()
        elif label == "ORG":
            return self.fake.company()
        elif label == "EMAIL":
            return self.fake.email()
        elif label == "PHONE":
            return self.fake.phone_number()
        elif label == "MONEY":
            return self.fake.pricetag()
        elif label == "IP":
            return self.fake.ipv4()
        elif label == "DATE":
            return self.fake.date()
        elif label == "TIME":
            return self.fake.time()
        elif label == "ADDRESS":
            return self.fake.address()
        else:
            return "[SYNTHETIC DATA]"

    def redact_presentation(self, ppt_path, redact_level, redaction_type):
        text, text_blocks = self.extract_text_from_presentation(ppt_path)
        redaction_labels = self.redaction_map.get(redact_level, [])

        sensitive_data = self.detect_sensitive_data(text, redaction_labels)

        prs = Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    original_text = shape.text
                    for entity, label in sensitive_data:
                        if entity in original_text:
                            if redaction_type == "blackout":
                                redacted = "█" * len(entity)
                            elif redaction_type == "blur":
                                redacted = "*" * len(entity)
                            elif redaction_type == "synthetic":
                                redacted = self.generate_synthetic_data(label)
                            else:
                                redacted = "[REDACTED]"
                            original_text = re.sub(re.escape(entity), redacted, original_text)
                    shape.text = original_text

        redacted_ppt_path = os.path.join(os.path.dirname(ppt_path), "redacted_" + os.path.basename(ppt_path))
        prs.save(redacted_ppt_path)
        return redacted_ppt_path